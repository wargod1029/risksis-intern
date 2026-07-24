"""
Generate a PowerPoint presentation from information.md
Standard AMD Workflow to Deploy AI Model on an FPGA
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── colours ──────────────────────────────────────────────
DARK_BG      = RGBColor(0x1B, 0x1B, 0x2F)   # deep navy
ACCENT_BLUE  = RGBColor(0x00, 0xBF, 0xFF)   # cyan accent
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x96)   # green accent
ACCENT_ORANGE= RGBColor(0xFF, 0xA5, 0x00)   # orange accent
ACCENT_RED   = RGBColor(0xFF, 0x55, 0x55)   # red accent
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCC, 0xCC, 0xCC)
TABLE_HDR_BG = RGBColor(0x2A, 0x3A, 0x5C)
TABLE_ROW_ALT= RGBColor(0x24, 0x2E, 0x4A)
TABLE_ROW_BG = RGBColor(0x1E, 0x26, 0x3E)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)


# ── helpers ──────────────────────────────────────────────
def dark_slide():
    """Return a blank slide with dark background."""
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    return slide


def add_bar(slide, left, top, width, height, color):
    """Add a coloured rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a simple text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_slide_number(slide, num, total):
    """Footer slide number."""
    add_textbox(slide, Inches(11.8), Inches(7.0), Inches(1.2), Inches(0.4),
                f"{num} / {total}", font_size=10, color=LIGHT_GRAY,
                alignment=PP_ALIGN.RIGHT)


def add_title_block(slide, title, subtitle=None, slide_num=None, total=None):
    """Consistent title bar at top of content slides."""
    add_bar(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.15), ACCENT_BLUE)
    add_textbox(slide, Inches(0.8), Inches(0.08), Inches(11.5), Inches(0.7),
                title, font_size=30, color=DARK_BG, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(0.65), Inches(11.5), Inches(0.4),
                    subtitle, font_size=14, color=DARK_BG, bold=False)
    if slide_num:
        add_slide_number(slide, slide_num, total)
    # thin accent line
    add_bar(slide, Inches(0), Inches(1.15), prs.slide_width, Inches(0.04), ACCENT_GREEN)


def add_bullet_frame(slide, left, top, width, height, bullets, font_size=16,
                     color=WHITE, spacing=Pt(6), font_name="Calibri"):
    """Add a text frame with bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
        p.level = 0
        p.bullet = True
    return txBox


def add_table(slide, left, top, col_widths, headers, rows, font_size=12):
    """Add a styled table. Returns table shape."""
    n_rows = 1 + len(rows)
    n_cols = len(headers)
    total_w = sum(col_widths)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, total_w,
                                         Inches(0.35 * n_rows))
    table = table_shape.table

    # set column widths
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HDR_BG
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER

    # data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ROW_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ROW_ALT
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.color.rgb = WHITE
                p.font.name = "Calibri"
                p.alignment = PP_ALIGN.CENTER

    return table_shape


TOTAL_SLIDES = 13


# ═══════════════════════════════════════════════════════════
# SLIDE  1 — TITLE
# ═══════════════════════════════════════════════════════════
slide = dark_slide()
add_bar(slide, Inches(0), Inches(1.8), Inches(0.12), Inches(3.8), ACCENT_BLUE)
add_textbox(slide, Inches(1.0), Inches(1.8), Inches(11.5), Inches(1.0),
            "AMD FPGA Workflow for",
            font_size=44, color=WHITE, bold=True)
add_textbox(slide, Inches(1.0), Inches(2.75), Inches(11.5), Inches(0.7),
            "AI Model Deployment on FPGA / SoC",
            font_size=36, color=ACCENT_BLUE, bold=False)
add_bar(slide, Inches(1.0), Inches(3.55), Inches(3.0), Inches(0.06), ACCENT_GREEN)
add_textbox(slide, Inches(1.0), Inches(3.9), Inches(11.5), Inches(0.5),
            "Vitis AI  •  Vitis HLS  •  Vivado  •  Vitis  —  The AMD FPGA AI Pipeline",
            font_size=18, color=LIGHT_GRAY)
add_textbox(slide, Inches(1.0), Inches(4.6), Inches(11.5), Inches(0.5),
            "July 2026", font_size=16, color=LIGHT_GRAY)
add_slide_number(slide, 1, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════
# SLIDE  2 — AMD ECOSYSTEM OVERVIEW
# ═══════════════════════════════════════════════════════════
slide = dark_slide()
add_title_block(slide, "AMD FPGA AI Ecosystem — Four Core Tools",
                "Target hardware: Zynq MPSoC / Zynq UltraScale+ / Versal ACAP", 2, TOTAL_SLIDES)

tools = [
    ("Vitis AI", "AI-model quantization,\ncompilation, and runtime\n(AI-specific pipeline)",
     "🧠", ACCENT_BLUE),
    ("Vitis HLS", "High-Level Synthesis\nC/C++ → RTL for custom\noperators DPU cannot handle",
     "⚙️", ACCENT_GREEN),
    ("Vivado", "Traditional FPGA design:\nRTL synthesis, place-and-route,\nIP integration, bitstream gen",
     "🔧", ACCENT_ORANGE),
    ("Vitis", "Unified software platform:\nhost app, PS–PL linking,\nplatform creation, debug",
     "💻", ACCENT_RED),
]

for i, (name, desc, icon, clr) in enumerate(tools):
    left = Inches(0.6 + i * 3.15)
    # card background
    add_bar(slide, left, Inches(1.7), Inches(2.85), Inches(3.8), RGBColor(0x1E, 0x26, 0x3E))
    # colored header
    add_bar(slide, left, Inches(1.7), Inches(2.85), Inches(0.7), clr)
    add_textbox(slide, left + Inches(0.1), Inches(1.75), Inches(2.65), Inches(0.6),
                f"{icon}  {name}", font_size=20, color=DARK_BG, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.15), Inches(2.6), Inches(2.55), Inches(2.7),
                desc, font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)

# bottom note
add_bar(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.7), RGBColor(0x1E, 0x30, 0x50))
add_textbox(slide, Inches(1.1), Inches(5.9), Inches(10.8), Inches(0.5),
            '💡  Zynq MPSoC / Versal ACAP combine ARM Processing System (PS) + FPGA Programmable Logic (PL)  '
            '— DPU sits in PL while Linux + VART run on PS.',
            font_size=14, color=ACCENT_GREEN)


# ═══════════════════════════════════════════════════════════
# SLIDE  3 — WORKFLOW OVERVIEW
# ═══════════════════════════════════════════════════════════
slide = dark_slide()
add_title_block(slide, "Vitis AI Workflow — End-to-End Pipeline",
                "From trained model to running DPU on FPGA", 3, TOTAL_SLIDES)

steps = [
    ("①  Export", "PyTorch /\nTensorFlow\n→ ONNX", ACCENT_BLUE),
    ("②  Quantize", "Vitis AI\nQuantizer\nfloat32 → INT8", ACCENT_BLUE),
    ("③  Compile", "Vitis AI\nCompiler\n→ .xmodel", ACCENT_GREEN),
    ("④  Hardware", "Vivado + Vitis\nDPU bitstream\n+ platform", ACCENT_ORANGE),
    ("⑤  Deploy", "VART / Vitis AI\nLibrary\nHost App", ACCENT_GREEN),
]

for i, (label, desc, clr) in enumerate(steps):
    left = Inches(0.5 + i * 2.55)
    add_bar(slide, left, Inches(1.8), Inches(2.2), Inches(1.1), clr)
    add_textbox(slide, left + Inches(0.05), Inches(1.85), Inches(2.1), Inches(0.35),
                label, font_size=16, color=DARK_BG, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.05), Inches(2.2), Inches(2.1), Inches(0.7),
                desc, font_size=14, color=WHITE, alignment=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        add_textbox(slide, left + Inches(2.22), Inches(2.0), Inches(0.35), Inches(0.5),
                    "▸", font_size=28, color=ACCENT_BLUE, bold=True)

# ASCII pipeline diagram
add_bar(slide, Inches(0.8), Inches(3.4), Inches(11.5), Inches(2.8), RGBColor(0x0D, 0x18, 0x2B))
code_text = (
    "PyTorch / TensorFlow (float32)\n"
    "        │\n"
    "        ▼\n"
    "   ①  Export to ONNX (or float model)\n"
    "        │\n"
    "        ▼\n"
    "   ②  Vitis AI Quantizer   →  INT8 model\n"
    "        │\n"
    "        ▼\n"
    "   ③  Vitis AI Compiler    →  .xmodel file\n"
    "        │\n"
    "        ▼\n"
    "   ④  Vivado + Vitis       →  DPU bitstream + platform\n"
    "        │\n"
    "        ▼\n"
    "   ⑤  VART / Vitis AI Library  →  host application (C++ / Python)"
)
add_textbox(slide, Inches(1.2), Inches(3.55), Inches(10.8), Inches(2.55),
            code_text, font_size=15, color=LIGHT_GRAY, font_name="Consolas")


# ═══════════════════════════════════════════════════════════
# SLIDE  4 — STEP 1 & 2: TRAIN/EXPORT + QUANTIZATION
# ═══════════════════════════════════════════════════════════
slide = dark_slide()
add_title_block(slide, "Step 1–2: Train, Export & Quantization",
                "From float32 model to INT8 deployment-ready model", 4, TOTAL_SLIDES)

# Step 1 box
add_bar(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.55), ACCENT_BLUE)
add_textbox(slide, Inches(0.95), Inches(1.63), Inches(5.2), Inches(0.5),
            "Step 1 — Train & Export", font_size=20, color=DARK_BG, bold=True)
add_bullet_frame(slide, Inches(0.8), Inches(2.3), Inches(5.5), Inches(2.0), [
    "Train model in PyTorch or TensorFlow (float32)",
    "Export to ONNX (Open Neural Network Exchange)",
    "Or pass float model directly if framework natively supported",
    "Vitis AI accepts ONNX + native TF/PyTorch formats",
], font_size=15, color=WHITE, spacing=Pt(8))

# Step 2 box
add_bar(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(0.55), ACCENT_GREEN)
add_textbox(slide, Inches(7.15), Inches(1.63), Inches(5.2), Inches(0.5),
            "Step 2 — Quantization (Vitis AI Quantizer)", font_size=20, color=DARK_BG, bold=True)
add_bullet_frame(slide, Inches(7.0), Inches(2.3), Inches(5.5), Inches(2.8), [
    "Convert weights & activations: float32 → INT8",
    "Post-Training Quantization (PTQ) with small calibration dataset",
    "100–1000 unlabeled samples → analyse activation distributions",
    "If accuracy loss unacceptable: Quantization-Aware Training (QAT)",
    "QAT tools: Brevitas (PyTorch) / QKeras / TF Model Optimization",
    "Fine-tune model with simulated quantization during training",
], font_size=15, color=WHITE, spacing=Pt(8))

# bottom callout
add_bar(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.7), RGBColor(0x1E, 0x30, 0x50))
add_textbox(slide, Inches(1.1), Inches(5.6), Inches(10.8), Inches(0.5),
            '🔑  INT8 quantization reduces model size ~4× and enables efficient DPU execution '
            'with minimal accuracy drop for most CNN architectures.',
            font_size=14, color=ACCENT_GREEN)


# ═══════════════════════════════════════════════════════════
# SLIDE  5 — STEP 3: COMPILATION
# ═══════════════════════════════════════════════════════════
slide = dark_slide()
add_title_block(slide, "Step 3 — Vitis AI Compiler",
                "Mapping the quantized model onto the DPU architecture", 5, TOTAL_SLIDES)

add_bullet_frame(slide, Inches(0.8), Inches(1.6), Inches(7.0), Inches(3.5), [
    "Feed quantized INT8 model into Vitis AI Compiler",
    "Compiler maps model graph onto target DPU architecture",
    "DPU configuration specified via arch.json file",
    "Performs operator fusion optimizations:",
    "  → BatchNorm folded into Convolution",
    "  → Activation layers merged with preceding ops",
    "Data-layout transformations for efficient DPU memory access",
    "Instruction scheduling across DPU cores (multi-core support)",
    "Output: .xmodel file — compiled instruction stream for DPU runtime",
], font_size=17, color=WHITE, spacing=Pt(8))

# compiler flow card
add_bar(slide, Inches(8.3), Inches(1.6), Inches(4.2), Inches(3.2), RGBColor(0x1E, 0x26, 0x3E))
add_textbox(slide, Inches(8.5), Inches(1.7), Inches(3.8), Inches(2.9),
            '📦  Compiler Input / Output\n\n'
            'INPUT:\n'
            '  • INT8 model (quantized)\n'
            '  • arch.json (DPU config)\n\n'
            'OUTPUT:\n'
            '  • .xmodel file\n\n'
            'The .xmodel contains:\n'
            '  • DPU instructions\n'
            '  • Weight data (INT8)\n'
            '  • Graph execution plan\n'
            '  • Multi-core scheduling',
            font_size=14, color=ACCENT_BLUE)

add_bar(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.7), RGBColor(0x1E, 0x30, 0x50))
add_textbox(slide, Inches(1.1), Inches(5.6), Inches(10.8), Inches(0.5),
            '⚡  For most standard CNNs (ResNet, YOLO, MobileNet, VGG, etc.) the DPU handles all operators directly '
            '— no custom HLS accelerators needed.',
            font_size=14, color=ACCENT_GREEN)


# ═══════════════════════════════════════════════════════════
# SLIDE  6 — STEP 4: HARDWARE PLATFORM (Vivado + Vitis)
# ═══════════════════════════════════════════════════════════
slide = dark_slide()
add_title_block(slide, "Step 4 — Hardware Platform (Vivado + Vitis)",
                "Creating the FPGA bitstream with the DPU", 6, TOTAL_SLIDES)

# Vivado section
add_bar(slide, Inches(0.8), Inches(1.6), Inches(5.7), Inches(0.5), ACCENT_ORANGE)
add_textbox(slide, Inches(0.95), Inches(1.62), Inches(5.4), Inches(0.45),
            "🔧  Vivado — FPGA Hardware Design", font_size=18, color=DARK_BG, bold=True)
add_bullet_frame(slide, Inches(0.8), Inches(2.2), Inches(5.7), Inches(3.0), [
    "Integrate DPU IP core into block design",
    "Configure DPU: #cores, BRAM/URAM allocation, conv types",
    "Add supporting IP: clocking, resets, AXI interconnects",
    "Add sensor interfaces: MIPI/CSI camera receivers, etc.",
    "Run synthesis → implementation (place-and-route)",
    "Generate bitstream (.bit) + hardware hand-off (.xsa)",
], font_size=15, color=WHITE, spacing=Pt(6))

# Vitis section
add_bar(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(0.5), ACCENT_RED)
add_textbox(slide, Inches(7.15), Inches(1.62), Inches(5.2), Inches(0.45),
            "💻  Vitis — Platform Creation", font_size=18, color=DARK_BG, bold=True)
add_bullet_frame(slide, Inches(7.0), Inches(2.2), Inches(5.5), Inches(3.0), [
    "Import .xsa from Vivado → create platform project",
    "Packages hardware design with software stack:",
    "  → FSBL (First Stage Boot Loader)",
    "  → PMU firmware (Platform Management Unit)",
    "  → Linux kernel + device tree",
    "Enables host application ↔ DPU communication",
    "Unified environment for system-level debug",
], font_size=15, color=WHITE, spacing=Pt(6))

# HLS callout
add_bar(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.0), RGBColor(0x1E, 0x30, 0x50))
add_textbox(slide, Inches(1.1), Inches(5.6), Inches(10.8), Inches(0.8),
            '⚙️  Vitis HLS (optional / advanced): For custom operators the DPU cannot natively support — write '
            'accelerator in C/C++, synthesize to RTL via Vitis HLS, integrate IP block alongside DPU in Vivado. '
            'Most standard CNNs do NOT need this.',
            font_size=14, color=ACCENT_GREEN)


# ═══════════════════════════════════════════════════════════
# SLIDE  7 — STEP 5: DEPLOYMENT & HOST APPLICATION
# ═══════════════════════════════════════════════════════════
slide = dark_slide()
add_title_block(slide, "Step 5 — Deployment & Host Application",
                "Flash bitstream + write host app with VART", 7, TOTAL_SLIDES)

# Flashing methods
add_textbox(slide, Inches(0.8), Inches(1.45), Inches(5.5), Inches(0.4),
            "Flashing the Bitstream", font_size=20, color=ACCENT_BLUE, bold=True)

col_w = [Inches(2.6), Inches(3.2)]
add_table(slide, Inches(0.8), Inches(1.95), col_w,
    ["Method", "Use Case"],
    [
        ["Vivado Hardware Manager (JTAG/USB)", "Development / debugging"],
        ["SD card / eMMC boot", "Standalone embedded (Zynq MPSoC, Kria)"],
        ["QSPI flash (program_flash)", "Production, non-volatile boot"],
        ["Ethernet / PCIe", "Data-centre accelerator (Alveo)"],
    ], font_size=13)

# Host application
add_textbox(slide, Inches(0.8), Inches(3.65), Inches(11.5), Inches(0.4),
            "Writing the Host Application — Vitis AI Runtime (VART)", font_size=20, color=ACCENT_BLUE, bold=True)

add_bullet_frame(slide, Inches(0.8), Inches(4.1), Inches(6.0), Inches(2.5), [
    "C++ or Python API to load .xmodel and run inference",
    "Submit inference jobs to DPU asynchronously",
    "Vitis AI Library: higher-level wrappers with pre-built pipelines",
    "Pre-processing: resize, normalize",
    "Post-processing: NMS (object detection), softmax (classification)",
    "Typical loop: capture → pre-process → DPU inference → post-process → act",
], font_size=14, color=WHITE, spacing=Pt(6))

# deployment architecture box
add_bar(slide, Inches(7.2), Inches(3.65), Inches(5.3), Inches(2.8), RGBColor(0x0D, 0x18, 0x2B))
arch_text = (
    "┌──────────────────────────────────────┐\n"
    "│  Processing System (ARM / Linux)     │\n"
    "│  ┌────────────────────────────────┐  │\n"
    "│  │  Host App (C++ / Python)       │  │\n"
    "│  │  VART / Vitis AI Library       │  │\n"
    "│  └───────────┬────────────────────┘  │\n"
    "│              │ AXI bus                │\n"
    "│  ┌───────────▼────────────────────┐  │\n"
    "│  │  Programmable Logic (FPGA)     │  │\n"
    "│  │  ┌──────┐  ┌──────┐  ┌──────┐ │  │\n"
    "│  │  │ DPU  │  │ MIPI │  │Custom│ │  │\n"
    "│  │  │ Core │  │ CSI  │  │ HLS  │ │  │\n"
    "│  │  └──────┘  └──────┘  └──────┘ │  │\n"
    "│  └────────────────────────────────┘  │\n"
    "└──────────────────────────────────────┘"
)
add_textbox(slide, Inches(7.35), Inches(3.75), Inches(5.0), Inches(2.6),
            arch_text, font_size=10, color=LIGHT_GRAY, font_name="Consolas")


# ═══════════════════════════════════════════════════════════
# SLIDE  8 — TOOL SUMMARY
# ═══════════════════════════════════════════════════════════
slide = dark_slide()
add_title_block(slide, "Tool Summary — End-to-End Workflow",
                "Complete tool chain from training to inference", 8, TOTAL_SLIDES)

col_w2 = [Inches(1.6), Inches(2.8), Inches(3.8), Inches(4.0)]
add_table(slide, Inches(0.4), Inches(1.6), col_w2,
    ["Step", "Tool", "Input", "Output"],
    [
        ["Train",     "PyTorch / TensorFlow",         "Dataset",                          "float32 model"],
        ["Export",    "torch.onnx / tf2onnx",         "float32 model",                    "ONNX"],
        ["Quantize",  "Vitis AI Quantizer",           "float32 model + calibration data", "INT8 model"],
        ["Compile",   "Vitis AI Compiler",            "INT8 model + arch.json",           ".xmodel"],
        ["Hardware",  "Vivado",                       "DPU IP + block design",            "bitstream (.bit) + .xsa"],
        ["Platform",  "Vitis",                        ".xsa",                             "platform (boot, device tree)"],
        ["Program",   "Vivado HW Manager / SD boot",  "bitstream",                        "running DPU on FPGA"],
        ["Host App",  "VART / Vitis AI Library",      ".xmodel + sensor data",            "inference results"],
    ], font_size=13)

add_bar(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(0.7), RGBColor(0x1E, 0x30, 0x50))
add_textbox(slide, Inches(1.1), Inches(5.3), Inches(10.8), Inches(0.5),
            '🐳  Environment: Vitis AI tools ship as Docker containers (CPU & GPU variants).  '
            'Vivado & Vitis are native Linux/Windows installs.  DPU IP core is included with Vitis AI.',
            font_size=14, color=ACCENT_GREEN)


# ═══════════════════════════════════════════════════════════
# SLIDE  9 — EXTERNAL RESOURCES
# ═══════════════════════════════════════════════════════════
slide = dark_slide()
add_title_block(slide, "External Resources & Learning Materials",
                "Key references for FPGA-based AI deployment", 9, TOTAL_SLIDES)

resources = [
    ("🎓", "Advanced FPGA Course: Verilog Based Robotics & Signal Processing",
     "YouTube playlist — comprehensive FPGA course covering Verilog for robotics and DSP",
     "https://www.youtube.com/playlist?list=PLfPwG72dAOx6gkqrrZVdr63EaizecO0Cg"),
    ("🔬", "hls4ml — Machine Learning to FPGA",
     "Translate ML models directly into synthesizable VHDL or Verilog",
     "https://fastmachinelearning.org/hls4ml/"),
    ("🔥", "PyTorch AI on FPGA — FINN Workflow Tutorial",
     "Step-by-step tutorial for deploying PyTorch models on FPGA using the FINN framework",
     "https://hugobrh.dev/posts/PY2FPGA/"),
    ("🐍", "Python to FPGA — Tutorial Snippets",
     "Practical examples and code snippets for the Python→FPGA workflow",
     "https://github.com/0BAB1/tutorial-snippets/tree/main/8%20Python%20to%20FPGA"),
    ("🤖", "Educational Platform for FPGA Accelerated AI in Robotics",
     "Open-source platform demonstrating FPGA-accelerated AI for robotics applications",
     "https://github.com/nhma20/FPGA_AI"),
    ("📖", "AMD Vitis AI 3.5 — Official Workflow Documentation",
     "Official AMD documentation for the complete model deployment workflow",
     "https://xilinx.github.io/Vitis-AI/3.5/html/docs/workflow-model-deployment.html"),
]

for i, (icon, title, desc, url) in enumerate(resources):
    top = Inches(1.5 + i * 0.95)
    add_bar(slide, Inches(0.6), top, Inches(0.07), Inches(0.7), ACCENT_BLUE)
    add_textbox(slide, Inches(0.85), top, Inches(11.5), Inches(0.35),
                f"{icon}  {title}", font_size=15, color=ACCENT_BLUE, bold=True)
    add_textbox(slide, Inches(0.85), top + Inches(0.32), Inches(8.5), Inches(0.3),
                desc, font_size=11, color=LIGHT_GRAY)
    add_textbox(slide, Inches(0.85), top + Inches(0.58), Inches(11.5), Inches(0.25),
                f"🔗  {url}", font_size=10, color=ACCENT_GREEN)


# ═══════════════════════════════════════════════════════════
# SLIDE 10 — AI MODULE COMPARISON (Part 1)
# ═══════════════════════════════════════════════════════════
slide = dark_slide()
add_title_block(slide, "AI Module Comparison — Edge Platforms (1/2)",
                "NVIDIA Jetson, Hailo, and x86 edge platforms", 10, TOTAL_SLIDES)

# card layout — 3 cards
cards_p1 = [
    ("NVIDIA Jetson\nAGX Orin",
     "Up to 275 TOPS (INT8)\n15W – 60W\n\nHigh-end robotics &\nautonomous navigation.\nIntegrates CPU, Ampere\nGPU, and DLAs.",
     ACCENT_GREEN),
    ("NVIDIA Jetson\nOrin Nano (Super)",
     "Up to 67 TOPS (INT8)\n7W – 25W\n\nEntry-level embedded AI\n& smart cameras.\n\"Super\" variant (Dec 2024)\noffers significant TOPS\nuplift over original.",
     ACCENT_GREEN),
    ("Raspberry Pi 5\n+ Hailo-8 AI HAT",
     "~13 to 26 TOPS\n10W – 15W\n\nLow-power field robots.\nHailo-8L = 13 TOPS\nHailo-8 = 26 TOPS\nSeparate Hailo accelerator.",
     ACCENT_BLUE),
]

for i, (name, desc, clr) in enumerate(cards_p1):
    left = Inches(0.45 + i * 4.2)
    add_bar(slide, left, Inches(1.5), Inches(3.85), Inches(4.3), RGBColor(0x1E, 0x26, 0x3E))
    add_bar(slide, left, Inches(1.5), Inches(3.85), Inches(0.65), clr)
    add_textbox(slide, left + Inches(0.1), Inches(1.52), Inches(3.65), Inches(0.6),
                name, font_size=17, color=DARK_BG, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.15), Inches(2.3), Inches(3.55), Inches(3.3),
                desc, font_size=13, color=WHITE)


# ═══════════════════════════════════════════════════════════
# SLIDE 11 — AI MODULE COMPARISON (Part 2)
# ═══════════════════════════════════════════════════════════
slide = dark_slide()
add_title_block(slide, "AI Module Comparison — Edge Platforms (2/2)",
                "Intel x86 edge, Axelera, EdgeCortix, SiMa.ai", 11, TOTAL_SLIDES)

cards_p2 = [
    ("ASUS NUC 14",
     "20–40+ TOPS (NPU+GPU)\n28W – 65+W (system)\n\nIndustrial edge / server.\nFull x86 + PCIe.\nMeteor Lake NPU ≈ 11 TOPS\nArrow Lake NPU ≈ 13 TOPS\nLunar Lake NPU ≈ 48 TOPS",
     ACCENT_ORANGE),
    ("Axelera AI\nMini PC + M.2 Max",
     "Up to 214 TOPS\n20W – 40W\n\nIn-memory computing (PIM)\n+ RISC-V architecture.\nUltra-low latency at high\nTOPS. Optimized for\ncomputer vision workloads.",
     ACCENT_RED),
    ("EdgeCortix\nSAKURA",
     "Up to 60 TOPS\nUnder 10W\n\nFIM (Fabric in Memory)\narchitecture. Targets\nlow-power smart city &\nvision applications.\nExtremely efficient.",
     ACCENT_BLUE),
    ("SiMa.ai\nMLSoC",
     "Up to 50+ TOPS\n~5–10W (board)\n\nPurpose-built for edge AI.\nIntegrated CV/vision pipeline\nwith software-programmable\nNPU. Power varies by\nworkload.",
     ACCENT_GREEN),
]

for i, (name, desc, clr) in enumerate(cards_p2):
    left = Inches(0.3 + i * 3.25)
    add_bar(slide, left, Inches(1.5), Inches(2.95), Inches(4.3), RGBColor(0x1E, 0x26, 0x3E))
    add_bar(slide, left, Inches(1.5), Inches(2.95), Inches(0.65), clr)
    add_textbox(slide, left + Inches(0.1), Inches(1.52), Inches(2.75), Inches(0.6),
                name, font_size=17, color=DARK_BG, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.1), Inches(2.3), Inches(2.75), Inches(3.3),
                desc, font_size=12, color=WHITE)


# ═══════════════════════════════════════════════════════════
# SLIDE 12 — FULL COMPARISON TABLE
# ═══════════════════════════════════════════════════════════
slide = dark_slide()
add_title_block(slide, "AI Module Comparison — Full Summary",
                "Side-by-side comparison of all evaluated edge AI platforms", 12, TOTAL_SLIDES)

col_w3 = [Inches(2.8), Inches(2.0), Inches(2.0), Inches(5.8)]
add_table(slide, Inches(0.4), Inches(1.5), col_w3,
    ["Platform / Device", "Peak AI (TOPS)", "Power", "Primary Use Case & Architecture"],
    [
        ["NVIDIA Jetson AGX Orin",           "Up to 275",   "15W – 60W",    "High-end robotics, autonomous nav. CPU + Ampere GPU + DLA."],
        ["NVIDIA Jetson Orin Nano (Super)",  "Up to 67",    "7W – 25W",     "Entry-level embedded AI & smart cameras. Super variant (Dec 2024) uplift."],
        ["Raspberry Pi 5 + Hailo-8 AI HAT",  "13 – 26",     "10W – 15W",    "Low-power field robots. Hailo-8L (13 TOPS) / Hailo-8 (26 TOPS) accelerator."],
        ["ASUS NUC 14",                      "20 – 40+",    "28W – 65+W",   "Industrial edge / server. Full x86 + PCIe. Gen-dependent NPU performance."],
        ["Axelera AI Mini PC + M.2 Max",     "Up to 214",   "20W – 40W",    "In-memory computing (PIM) + RISC-V. Ultra-low latency at high TOPS."],
        ["EdgeCortix SAKURA",               "Up to 60",    "Under 10W",    "FIM architecture. Low-power smart city & vision."],
        ["SiMa.ai MLSoC",                    "Up to 50+",   "~5 – 10W",     "Purpose-built edge AI. Integrated CV pipeline + programmable NPU."],
    ], font_size=12)


# ═══════════════════════════════════════════════════════════
# SLIDE 13 — KEY TAKEAWAYS
# ═══════════════════════════════════════════════════════════
slide = dark_slide()
add_title_block(slide, "Key Takeaways — AMD FPGA AI Deployment",
                "Summary of the Vitis AI workflow", 13, TOTAL_SLIDES)

takeaways = [
    ("🧠", "Complete Pipeline",
     "AMD provides a full end-to-end tool chain: PyTorch/TF → ONNX → Quantize → Compile → "
     "Hardware → Deploy. Each step is handled by a dedicated tool in the ecosystem.",
     ACCENT_BLUE),
    ("⚡", "DPU Acceleration",
     "The Deep Learning Processing Unit (DPU) is a hardened IP core optimized for CNN inference. "
     "INT8 quantization + operator fusion + multi-core scheduling deliver efficient FPGA acceleration.",
     ACCENT_GREEN),
    ("🔧", "Hardware Flexibility",
     "Vivado + Vitis HLS allow custom hardware accelerators for non-standard operators. "
     "The FPGA can be tailored to the exact model architecture — not limited to fixed silicon.",
     ACCENT_ORANGE),
    ("🐳", "Docker-Based Tooling",
     "Vitis AI tools run in Docker containers (CPU/GPU variants) for easy setup. "
     "Vivado & Vitis are native installs for synthesis and platform creation.",
     ACCENT_BLUE),
    ("🌐", "Rich Ecosystem",
     "Multiple external resources: hls4ml, FINN, educational FPGA-AI platforms, and "
     "official AMD documentation. Active community and growing adoption in edge AI.",
     ACCENT_GREEN),
]

for i, (icon, title, desc, clr) in enumerate(takeaways):
    top = Inches(1.5 + i * 1.15)
    add_bar(slide, Inches(0.6), top, Inches(0.08), Inches(0.9), clr)
    add_textbox(slide, Inches(0.9), top, Inches(11.5), Inches(0.4),
                f"{icon}  {title}", font_size=20, color=clr, bold=True)
    add_textbox(slide, Inches(0.9), top + Inches(0.42), Inches(11.5), Inches(0.65),
                desc, font_size=14, color=WHITE)


# ── save ──────────────────────────────────────────────────
out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                        "AMD_FPGA_AI_Model_Deployment.pptx")
prs.save(out_path)
print(f"Saved to: {out_path}")
print(f"Slides: {len(prs.slides)}")
