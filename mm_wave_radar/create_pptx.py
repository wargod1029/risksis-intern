"""
Generate a PowerPoint presentation from information.md
mmWave Radar for Live Animal Detection in Passenger Screening
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── colours ──────────────────────────────────────────────
DARK_BG      = RGBColor(0x1B, 0x1B, 0x2F)   # deep navy
ACCENT_BLUE  = RGBColor(0x00, 0xBF, 0xFF)   # cyan accent
ACCENT_GREEN = RGBColor(0x00, 0xE6, 0x96)   # green accent
ACCENT_RED   = RGBColor(0xFF, 0x55, 0x55)   # red accent
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCC, 0xCC, 0xCC)
TABLE_HDR_BG = RGBColor(0x2A, 0x3A, 0x5C)
TABLE_ROW_ALT= RGBColor(0x24, 0x2E, 0x4A)
TABLE_ROW_BG = RGBColor(0x1E, 0x26, 0x3E)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)


# ── helpers ──────────────────────────────────────────────
def dark_slide():
    """Return a blank slide with dark background."""
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG
    return slide


def add_bar(slide, left, top, width, height, color):
    """Add a coloured rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a simple text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_slide_number(slide, num, total):
    """Footer slide number."""
    add_textbox(slide, Inches(11.8), Inches(7.0), Inches(1.2), Inches(0.4),
                f"{num} / {total}", font_size=10, color=LIGHT_GRAY,
                alignment=PP_ALIGN.RIGHT)


def add_title_block(slide, title, subtitle=None, slide_num=None, total=None):
    """Consistent title bar at top of content slides."""
    add_bar(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.15), ACCENT_BLUE)
    add_textbox(slide, Inches(0.8), Inches(0.08), Inches(11.5), Inches(0.7),
                title, font_size=30, color=DARK_BG, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(0.65), Inches(11.5), Inches(0.4),
                    subtitle, font_size=14, color=DARK_BG, bold=False)
    if slide_num:
        add_slide_number(slide, slide_num, total)
    # thin accent line
    add_bar(slide, Inches(0), Inches(1.15), prs.slide_width, Inches(0.04), ACCENT_GREEN)


def add_bullet_frame(slide, left, top, width, height, bullets, font_size=16,
                     color=WHITE, spacing=Pt(6), font_name="Calibri"):
    """Add a text frame with bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
        p.level = 0
        # bullet character
        p.bullet = True
    return txBox


def add_table(slide, left, top, col_widths, headers, rows, font_size=12):
    """Add a styled table. Returns table shape."""
    n_rows = 1 + len(rows)
    n_cols = len(headers)
    total_w = sum(col_widths)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, total_w,
                                         Inches(0.35 * n_rows))
    table = table_shape.table

    # set column widths
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HDR_BG
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER

    # data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ROW_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ROW_ALT
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.color.rgb = WHITE
                p.font.name = "Calibri"
                p.alignment = PP_ALIGN.CENTER

    return table_shape


TOTAL_SLIDES = 16


# ═══════════════════════════════════════════════════════════
# SLIDE  1 — TITLE
# ═══════════════════════════════════════════════════════════
slide = dark_slide()
add_bar(slide, Inches(0), Inches(1.8), Inches(0.12), Inches(3.8), ACCENT_BLUE)
add_textbox(slide, Inches(1.0), Inches(1.8), Inches(11.5), Inches(1.0),
            "mmWave Radar for Live Animal Detection",
            font_size=44, color=WHITE, bold=True)
add_textbox(slide, Inches(1.0), Inches(2.75), Inches(11.5), Inches(0.7),
            "in Passenger Screening",
            font_size=36, color=ACCENT_BLUE, bold=False)
add_bar(slide, Inches(1.0), Inches(3.55), Inches(3.0), Inches(0.06), ACCENT_GREEN)
add_textbox(slide, Inches(1.0), Inches(3.9), Inches(11.5), Inches(0.5),
            "Enhancing Biosecurity at Clearance Checkpoints with Advanced Radar Sensing",
            font_size=18, color=LIGHT_GRAY)
add_textbox(slide, Inches(1.0), Inches(4.6), Inches(11.5), Inches(0.5),
            "July 2026", font_size=16, color=LIGHT_GRAY)
add_slide_number(slide, 1, TOTAL_SLIDES)


# ═══════════════════════════════════════════════════════════
# SLIDE  2 — PROBLEM STATEMENT
# ═══════════════════════════════════════════════════════════
slide = dark_slide()
add_title_block(slide, "Problem Statement", "Illegal animal importation screening", 2, TOTAL_SLIDES)

add_bullet_frame(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(2.8), [
    "Government seeking advanced technology to enhance passenger screening at clearance halls",
    "Assist frontline officers in identifying and intercepting illegal animal importation",
    "Animals concealed within baggage, backpacks, or luggage",
    "Quick and accurate detection of live animals required",
    "Focus on warm-blooded animals — cold-blooded detection highly preferred",
], font_size=20, color=WHITE, spacing=Pt(12))

# summary box
add_bar(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(1.0), RGBColor(0x1E, 0x30, 0x50))
add_textbox(slide, Inches(1.1), Inches(5.15), Inches(10.8), Inches(0.7),
            '💡  Goal: A system that detects live animals through baggage non-invasively, '
            'with high throughput and no ionizing radiation.',
            font_size=16, color=ACCENT_GREEN)


# ═══════════════════════════════════════════════════════════
# SLIDE  3 — WHY mmWave RADAR
# ═══════════════════════════════════════════════════════════
slide = dark_slide()
add_title_block(slide, "Why mmWave Radar?", "Core technology differentiators", 3, TOTAL_SLIDES)

add_bullet_frame(slide, Inches(0.8), Inches(1.6), Inches(5.8), Inches(3.5), [
    "Electromagnetic waves penetrate non-metallic materials",
    "  → Clothing, backpacks, plastic containers, cardboard",
    "MIMO radar creates high-resolution 3D images",
    "Key advantage: 3D Doppler Imaging",
    "Detects micro-movements even when animal is still:",
    "  → Breathing (~0.1–3 Hz, 0.1–1 mm displacement)",
    "  → Heartbeats (~1–3 Hz, ~0.1–0.5 mm)",
    "  → Involuntary twitches",
], font_size=16, color=WHITE, spacing=Pt(6))

# right-side callout box
add_bar(slide, Inches(7.2), Inches(1.6), Inches(5.3), Inches(2.5), RGBColor(0x1E, 0x30, 0x50))
add_textbox(slide, Inches(7.5), Inches(1.75), Inches(4.8), Inches(2.2),
            '🔑  Core Insight\n\nEven a perfectly still animal cannot hide — '
            'breathing and heartbeat produce detectable phase modulation in '
            'the radar return. This is a physiological inevitability for all '
            'living vertebrates.',
            font_size=15, color=ACCENT_BLUE)

# technology comparison table
col_w = [Inches(1.8), Inches(1.8), Inches(1.8), Inches(1.8), Inches(2.2), Inches(1.8)]
add_table(slide, Inches(0.8), Inches(4.5), col_w,
    ["Technology", "Detects Life?", "Penetrates Bags?", "3D Imaging?", "Ionizing Radiation?", "Throughput"],
    [
        ["X-ray",           "No",                 "Yes",              "No (2D only)", "Yes",              "High"],
        ["Thermal IR",      "Indirectly (heat)",  "No (blocked)",     "No",           "No",               "Medium"],
        ["mmWave Radar",    "Yes (Doppler)",      "Yes",              "Yes (MIMO)",   "No",               "High"],
        ["Manual Inspection","Yes",               "Yes (opening)",    "N/A",          "No",               "Low"],
    ], font_size=11)


# ═══════════════════════════════════════════════════════════
# SLIDE  4 — SENSOR OVERVIEW
# ═══════════════════════════════════════════════════════════
slide = dark_slide()
add_title_block(slide, "Radar Sensor Candidates", "Four sensor options evaluated", 4, TOTAL_SLIDES)

cards = [
    ("TI IWR6843ISK",      "60 GHz | 3TX/4RX\n12 virtual antennas\nVital sign demo built-in\n~$299–$799 USD", ACCENT_BLUE),
    ("TI IWR1843BOOST",    "77 GHz | 3TX/4RX\n12 virtual antennas\nOn-chip DSP + MCU\n~$349–$849 USD", ACCENT_GREEN),
    ("TI IWRL6432BOOST",   "60 GHz | 2TX/3RX\n6 virtual antennas\nLow power, compact\n~$199–$699 USD", ACCENT_BLUE),
    ("Infineon BGT60TR13C","60 GHz | 1TX/3RX\n3 virtual channels\n5.5 GHz bandwidth\n~€250 EUR", ACCENT_GREEN),
]

for idx, (name, desc, color) in enumerate(cards):
    left = Inches(0.6 + idx * 3.1)
    add_bar(slide, left, Inches(1.5), Inches(2.8), Inches(0.6), color)
    add_textbox(slide, left + Inches(0.1), Inches(1.52), Inches(2.6), Inches(0.55),
                name, font_size=15, color=DARK_BG, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.15), Inches(2.3), Inches(2.5), Inches(2.8),
                desc, font_size=13, color=WHITE)


# ═══════════════════════════════════════════════════════════
# SLIDE  5 — TI IWR6843ISK  (Detail)
# ═══════════════════════════════════════════════════════════
slide = dark_slide()
add_title_block(slide, "TI IWR6843ISK — 60 GHz mmWave Radar Sensor",
                "Recommended starting point  |  ti.com/tool/IWR6843ISK", 5, TOTAL_SLIDES)

specs = [
    ("Frequency Band",      "60–64 GHz"),
    ("TX / RX Channels",    "3 Transmitters, 4 Receivers (12 virtual antennas)"),
    ("Bandwidth",           "Up to 4 GHz  →  range resolution ~3.75 cm"),
    ("Detection Range",     "Up to 50 m (people); 0.3–2 m for vital signs"),
    ("Velocity Resolution", "~0.1 mm/s Doppler micro-motion sensitivity"),
    ("Field of View",       "±60° Azimuth, ±60° Elevation"),
    ("Output Interface",    "USB (via DCA1000EVM) or direct SPI"),
    ("SDK / Processing",    "TI mmWave SDK, mmWave Studio, MATLAB/Python"),
    ("Vital Sign Demo",     "Built-in: breathing rate & heartbeat rate detection"),
    ("3D Imaging",          "Yes — MIMO virtual array + 3D point cloud + Doppler"),
    ("Dimensions",          "63.5 mm × 44.5 mm (standalone PCB)"),
    ("Price",               "~$299 USD (ISK) / ~$799 USD (with DCA1000EVM)"),
]

col_w2 = [Inches(3.5), Inches(8.0)]
rows = [[a, b] for a, b in specs]
tbl = add_table(slide, Inches(0.8), Inches(1.6), col_w2,
                ["Feature", "Specification"], rows, font_size=13)

# key strength callout
add_bar(slide, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.7), RGBColor(0x1E, 0x30, 0x50))
add_textbox(slide, Inches(1.1), Inches(6.48), Inches(10.8), Inches(0.55),
            '⭐  Most mature mmWave platform — TI provides ready-to-run lab demos for breathing/heartbeat detection. '
            'DCA1000EVM required for raw ADC data capture and custom 3D/Doppler processing.',
            font_size=14, color=ACCENT_GREEN)


# ═══════════════════════════════════════════════════════════
# SLIDE  6 — TI IWR1843BOOST
# ═══════════════════════════════════════════════════════════
slide = dark_slide()
add_title_block(slide, "TI IWR1843BOOST — 77 GHz Industrial Radar Sensor",
                "Higher frequency, on-chip processing  |  ti.com/tool/IWR1843BOOST", 6, TOTAL_SLIDES)

specs2 = [
    ("Frequency Band",      "76–81 GHz"),
    ("TX / RX Channels",    "3 TX, 4 RX (12 virtual antennas)"),
    ("Bandwidth",           "Up to 4 GHz → range resolution ~3.75 cm"),
    ("Detection Range",     "Up to 100 m; 0.3–2 m for micro-motion"),
    ("Velocity Resolution", "Sub-mm/s Doppler sensitivity"),
    ("Field of View",       "±60° Azimuth, ±60° Elevation"),
    ("On-Chip DSP",         "TI C674x DSP for real-time FFT & CFAR processing"),
    ("On-Chip MCU",         "ARM Cortex-R4F for tracking & classification"),
    ("3D Imaging",          "Yes — MIMO virtual array, 3D point cloud generation"),
    ("Dimensions",          "79 mm × 63.5 mm"),
    ("Price",               "~$349 USD (BOOST) / ~$849 USD (with DCA1000EVM)"),
]

add_table(slide, Inches(0.8), Inches(1.6), col_w2,
          ["Feature", "Specification"], [[a, b] for a, b in specs2], font_size=13)

add_bar(slide, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.7), RGBColor(0x1E, 0x30, 0x50))
add_textbox(slide, Inches(1.1), Inches(6.48), Inches(10.8), Inches(0.55),
            '🔧  Higher frequency (77 GHz) gives slightly better resolution + smaller antennas. '
            'On-chip DSP + MCU enables real-time embedded processing at the checkpoint — no external host required.',
            font_size=14, color=ACCENT_GREEN)


# ═══════════════════════════════════════════════════════════
# SLIDE  7 — TI IWRL6432BOOST
# ═══════════════════════════════════════════════════════════
slide = dark_slide()
add_title_block(slide, "TI IWRL6432BOOST — 57–64 GHz Low-Power Sensor",
                "Compact & efficient  |  ti.com/tool/IWRL6432BOOST", 7, TOTAL_SLIDES)

specs3 = [
    ("Frequency Band",   "57–64 GHz"),
    ("TX / RX Channels", "2 TX, 3 RX (6 virtual antennas)"),
    ("Bandwidth",        "Up to 4 GHz"),
    ("Detection Range",  "Up to 20 m; 0.3–2 m for micro-motion"),
    ("FOV",              "±60° Azimuth, ±50° Elevation"),
    ("Power",            "~2.5 mW deep sleep / ~1.5 W active"),
    ("Dimensions",       "43 mm × 25 mm"),
    ("Price",            "~$199 USD (BOOST only)"),
]

add_table(slide, Inches(0.8), Inches(1.6), col_w2,
          ["Feature", "Specification"], [[a, b] for a, b in specs3], font_size=13)

add_bar(slide, Inches(0.8), Inches(5.3), Inches(5.5), Inches(1.5), RGBColor(0x1E, 0x30, 0x50))
add_textbox(slide, Inches(1.1), Inches(5.4), Inches(5.0), Inches(1.3),
            '✅  Advantages\n• Lowest power consumption\n• Smallest form factor\n• Best for battery-powered /'
            '\n   space-constrained deployments\n• Lowest cost entry point',
            font_size=14, color=ACCENT_GREEN)

add_bar(slide, Inches(6.9), Inches(5.3), Inches(5.5), Inches(1.5), RGBColor(0x3A, 0x20, 0x20))
add_textbox(slide, Inches(7.2), Inches(5.4), Inches(5.0), Inches(1.3),
            '⚠️  Limitations\n• 2TX/3RX — only 6 virtual channels\n• Reduced 3D angular resolution vs.\n   '
            'IWR6843/IWR1843 (12 virtual channels)\n• May not match 3D imaging quality',
            font_size=14, color=ACCENT_RED)


# ═══════════════════════════════════════════════════════════
# SLIDE  8 — INFINEON BGT60TR13C
# ═══════════════════════════════════════════════════════════
slide = dark_slide()
add_title_block(slide, "Infineon DEMO BGT60TR13C — 60 GHz Radar Sensor",
                "Compact sensor, high bandwidth  |  infineon.com", 8, TOTAL_SLIDES)

specs4 = [
    ("Frequency Band",   "57–63 GHz (V-band)"),
    ("TX / RX Channels", "1 TX, 3 RX (3 virtual channels)"),
    ("Bandwidth",        "Up to 5.5 GHz → range resolution ~2.7 cm"),
    ("Detection Range",  "Up to 15 m; 0.2–1 m for vital signs"),
    ("FOV",              "±45° Azimuth, ±45° Elevation"),
    ("SDK",              "Infineon RDK, Radar Fusion GUI, Python SDK"),
    ("Vital Sign Demo",  "Yes — dedicated vital sensing reference app"),
    ("Sensor Size",      "25 mm × 25 mm"),
    ("Price",            "~€250 EUR"),
]

add_table(slide, Inches(0.8), Inches(1.6), col_w2,
          ["Feature", "Specification"], [[a, b] for a, b in specs4], font_size=13)

add_bar(slide, Inches(0.8), Inches(5.8), Inches(5.5), Inches(1.3), RGBColor(0x1E, 0x30, 0x50))
add_textbox(slide, Inches(1.1), Inches(5.9), Inches(5.0), Inches(1.1),
            '✅  Advantages\n• Highest bandwidth (5.5 GHz)\n   → finest range resolution (~2.7 cm)\n• '
            'Ultra-compact sensor (25×25 mm)\n• Can cascade multiple for custom MIMO',
            font_size=14, color=ACCENT_GREEN)

add_bar(slide, Inches(6.9), Inches(5.8), Inches(5.5), Inches(1.3), RGBColor(0x3A, 0x20, 0x20))
add_textbox(slide, Inches(7.2), Inches(5.9), Inches(5.0), Inches(1.1),
            '⚠️  Limitations\n• Single TX channel\n   → no native MIMO imaging\n• Best for single-point vital sign\n   '
            'detection, not full 3D imaging\n• Requires multiple synced boards\n   for MIMO',
            font_size=14, color=ACCENT_RED)


# ═══════════════════════════════════════════════════════════
# SLIDE  9 — RADAR COMPARISON SUMMARY
# ═══════════════════════════════════════════════════════════
slide = dark_slide()
add_title_block(slide, "Radar Comparison Summary", "Side-by-side sensor comparison", 9, TOTAL_SLIDES)

col_w3 = [Inches(2.4), Inches(1.4), Inches(1.5), Inches(1.7), Inches(2.2), Inches(2.0), Inches(1.6)]
add_table(slide, Inches(0.4), Inches(1.6), col_w3,
    ["Sensor", "Frequency", "TX/RX", "Max BW", "3D Imaging", "Vital Sign Demo", "Price (USD)"],
    [
        ["TI IWR6843ISK",     "60 GHz", "3/4 (12 virt.)", "4 GHz", "Yes (moderate)", "Yes",              "$299–$799"],
        ["TI IWR1843BOOST",   "77 GHz", "3/4 (12 virt.)", "4 GHz", "Yes (moderate)", "Via custom DSP",   "$349–$849"],
        ["TI IWRL6432BOOST",  "60 GHz", "2/3 (6 virt.)",  "4 GHz", "Limited",        "Via custom DSP",   "$199–$699"],
        ["Infineon BGT60TR13C","60 GHz","1/3 (3 virt.)",  "5.5 GHz","No (single TX)", "Yes",              "~€250"],
    ], font_size=13)

# recommendation
add_bar(slide, Inches(0.8), Inches(4.3), Inches(11.5), Inches(1.8), RGBColor(0x1E, 0x30, 0x50))
add_textbox(slide, Inches(1.1), Inches(4.4), Inches(10.8), Inches(1.6),
            '🏆  Recommendation\n\n'
            'Start with TI IWR6843ISK — lowest barrier to entry, built-in vital sign demo, '
            'mature SDK, and proven 3D MIMO imaging capability at the lowest price point.\n'
            'For higher-resolution 3D imaging, consider cascading multiple IWR6843 units or IWR1843BOOST.',
            font_size=15, color=ACCENT_GREEN)


# ═══════════════════════════════════════════════════════════
# SLIDE 10 — DETECTION MECHANISM  (Warm-blooded)
# ═══════════════════════════════════════════════════════════
slide = dark_slide()
add_title_block(slide, "Detection Mechanism — Warm-Blooded Animals",
                "Birds & Mammals", 10, TOTAL_SLIDES)

add_bullet_frame(slide, Inches(0.8), Inches(1.6), Inches(7.5), Inches(3.0), [
    "Strong, rhythmic Doppler signatures from:",
    "   → Breathing: 0.1–1 mm chest wall displacement, 0.1–3 Hz",
    "   → Heartbeat: ~0.1–0.5 mm displacement, 1–3 Hz",
    "Body temperature (~37°C) vs. ambient creates dielectric contrast",
    "   → Enhances radar reflectivity at tissue boundaries",
    "",
    "Detection Confidence: HIGH ✅",
    "TI's vital-sign demos prove reliable extraction at 1–2 m",
    "range through clothing and non-metallic barriers.",
], font_size=17, color=WHITE, spacing=Pt(8))

# icon / visual indicator
add_bar(slide, Inches(9.0), Inches(1.6), Inches(3.5), Inches(3.5), RGBColor(0x0A, 0x3D, 0x0A))
add_textbox(slide, Inches(9.2), Inches(2.5), Inches(3.1), Inches(2.0),
            '🫁  Breathing\n    0.1–1 mm\n    0.1–3 Hz\n\n'
            '❤️  Heartbeat\n    ~0.1–0.5 mm\n    1–3 Hz\n\n'
            '🔥  Body Heat\n    Dielectric\n    Contrast',
            font_size=16, color=ACCENT_GREEN)


# ═══════════════════════════════════════════════════════════
# SLIDE 11 — DETECTION MECHANISM  (Cold-blooded)
# ═══════════════════════════════════════════════════════════
slide = dark_slide()
add_title_block(slide, "Detection Mechanism — Cold-Blooded Animals",
                "Reptiles & Amphibians", 11, TOTAL_SLIDES)

add_bullet_frame(slide, Inches(0.8), Inches(1.6), Inches(7.5), Inches(3.0), [
    "Weaker but present respiratory micro-motions:",
    "   → Periodic throat / body wall movements",
    "   → Lung ventilation at 0.05–0.5 Hz",
    "Lower body temperature reduces dielectric contrast but:",
    "   → Tissue structures (lungs, organs, skin) still reflect",
    "",
    "Detection Confidence: MEDIUM–HIGH ⚠️",
    "Dependent on animal size, activity level, and concealment.",
    "Longer integration (5–15 s) increases SNR significantly.",
], font_size=17, color=WHITE, spacing=Pt(8))

add_bar(slide, Inches(9.0), Inches(1.6), Inches(3.5), Inches(3.5), RGBColor(0x3D, 0x2A, 0x0A))
add_textbox(slide, Inches(9.2), Inches(2.5), Inches(3.1), Inches(2.0),
            '🦎  Reptiles\n    Throat / body\n    movements\n    0.05–0.5 Hz\n\n'
            '🐍  Amphibians\n    Lung ventilation\n    Periodic motion\n\n'
            '⏱️  Longer Integration\n    5–15 seconds\n    → Higher SNR',
            font_size=16, color=ACCENT_BLUE)


# ═══════════════════════════════════════════════════════════
# SLIDE 12 — KEY ADVANTAGE
# ═══════════════════════════════════════════════════════════
slide = dark_slide()
add_title_block(slide, "Key Advantage Over Existing Technologies",
                "Doppler = physiological inevitability", 12, TOTAL_SLIDES)

add_bar(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(3.0), RGBColor(0x1E, 0x30, 0x50))

add_textbox(slide, Inches(1.2), Inches(2.0), Inches(10.8), Inches(2.5),
            'Unlike thermal cameras (which require the animal\'s body heat to reach the sensor '
            'with an unobstructed view) or X-ray (which sees only density), '
            'mmWave Doppler detection relies on physical motion — a physiological inevitability '
            'for all living vertebrates.\n\n'
            'Even at rest, breathing and heartbeat produce detectable phase modulation in the '
            'radar return.\n\n'
            'This cannot be masked by the animal — it is involuntary and continuous.',
            font_size=19, color=WHITE)

# 3 comparison boxes
comps = [
    ("🌡️  Thermal IR", "Blocked by fabric\nNo 3D imaging\nIndirect detection only", ACCENT_RED),
    ("☢️  X-ray", "Ionizing radiation\n2D only — no life detection\nSees density, not motion", ACCENT_RED),
    ("📡  mmWave Radar", "Penetrates non-metallic\n3D + Doppler detection\nDetects life itself", ACCENT_GREEN),
]
for i, (title, desc, clr) in enumerate(comps):
    left = Inches(0.8 + i * 4.1)
    add_bar(slide, left, Inches(5.3), Inches(3.7), Inches(1.8), clr)
    add_textbox(slide, left + Inches(0.15), Inches(5.4), Inches(3.4), Inches(0.5),
                title, font_size=17, color=DARK_BG, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.15), Inches(5.9), Inches(3.4), Inches(1.1),
                desc, font_size=14, color=DARK_BG, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 13 — SYSTEM ARCHITECTURE
# ═══════════════════════════════════════════════════════════
slide = dark_slide()
add_title_block(slide, "System Architecture Concept", "Clearance checkpoint deployment", 13, TOTAL_SLIDES)

# architecture flow boxes
steps = [
    ("Bag on\nConveyor", ACCENT_BLUE),
    ("MIMO mmWave\nRadar Array\n(e.g. cascaded\nIWR6843)", ACCENT_BLUE),
    ("Real-Time\nProcessor\n• Static 3D\n• Doppler\n• AI/ML", ACCENT_GREEN),
    ("Operator\nDisplay\n🚦 Green/Amber/Red", ACCENT_GREEN),
]

for i, (label, clr) in enumerate(steps):
    left = Inches(0.8 + i * 3.2)
    add_bar(slide, left, Inches(2.5), Inches(2.7), Inches(1.6), clr)
    add_textbox(slide, left + Inches(0.1), Inches(2.6), Inches(2.5), Inches(1.4),
                label, font_size=17, color=DARK_BG, bold=True, alignment=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        add_textbox(slide, left + Inches(2.72), Inches(2.9), Inches(0.5), Inches(0.5),
                    "▶", font_size=28, color=ACCENT_BLUE, bold=True)

# signal processing chain
add_textbox(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.5),
            "Signal Processing Chain", font_size=20, color=ACCENT_BLUE, bold=True)

chain = [
    ("1", "Static 3D\nReconstruction", "FMCW range processing\n+ MIMO angle estimation\n→ volumetric voxel map"),
    ("2", "Doppler Micro-Motion\nExtraction", "Coherent phase analysis\nSub-mm displacements\nPhysiological frequencies"),
    ("3", "AI/ML\nFusion", "Deep classifier on\nlabeled radar signatures\nShape + Doppler fusion"),
    ("4", "Alert\nOutput", "🟢 Green: Clear\n🟠 Amber: Ambiguous\n🔴 Red: Live animal"),
]

for i, (num, title, desc) in enumerate(chain):
    left = Inches(0.6 + i * 3.15)
    add_bar(slide, left, Inches(5.1), Inches(2.8), Inches(0.5), ACCENT_BLUE)
    add_textbox(slide, left + Inches(0.1), Inches(5.12), Inches(2.6), Inches(0.45),
                f"{num}.  {title}", font_size=13, color=DARK_BG, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.05), Inches(5.7), Inches(2.7), Inches(1.5),
                desc, font_size=12, color=WHITE, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 14 — TECHNICAL RISKS & MITIGATIONS
# ═══════════════════════════════════════════════════════════
slide = dark_slide()
add_title_block(slide, "Technical Risks & Mitigations", "", 14, TOTAL_SLIDES)

col_w4 = [Inches(4.0), Inches(7.5)]
add_table(slide, Inches(0.8), Inches(1.6), col_w4,
    ["Risk", "Mitigation"],
    [
        ["Low SNR for cold-blooded animals",
         "Longer integration (5–15 s), higher-gain antenna arrays, matched filtering for respiratory waveforms"],
        ["Metallic baggage lining (full reflection)",
         "Secondary X-ray check + flag metallic contents — detection impossible through foil-lined bags"],
        ["Conveyor motion clutter",
         "Stationary radar geometry + Doppler clutter filtering removes constant-velocity belt motion"],
        ["False positives (vibrating electronics, loose items)",
         "Spectral pattern analysis + ML classifier to distinguish physiological (variable) from mechanical (single-frequency) motion"],
        ["Throughput constraints",
         "Fast FMCW chirp rates + parallel beamforming → < 5 s per bag dwell time"],
    ], font_size=14)


# ═══════════════════════════════════════════════════════════
# SLIDE 15 — FEASIBILITY
# ═══════════════════════════════════════════════════════════
slide = dark_slide()
add_title_block(slide, "Feasibility & Supporting Research", "", 15, TOTAL_SLIDES)

# 3 research pillars
research = [
    ("🏥  Remote Vital Sign Monitoring",
     "Multiple peer-reviewed studies demonstrate mmWave detection of human heartbeat and respiration through clothing and blankets at 1–3 m (60–77 GHz FMCW). Same principles apply to animals. TI provides a production-ready vital signs demo with IWR6843."),
    ("🛂  Security Screening (Deployed)",
     "MIMO mmWave body scanners for concealed object detection already deployed at airports (e.g., Rohde & Schwarz QPS systems at 70–80 GHz). Bag-screening adaptation is a straightforward extension."),
    ("🛟  Search & Rescue Life Detection",
     "Lower-frequency Doppler radar is field-proven for detecting breathing survivors under rubble. mmWave brings higher spatial resolution appropriate for baggage-scale screening."),
]

for i, (title, desc) in enumerate(research):
    top = Inches(1.6 + i * 1.9)
    add_bar(slide, Inches(0.8), top, Inches(0.08), Inches(1.5), ACCENT_BLUE)
    add_textbox(slide, Inches(1.1), top, Inches(11.5), Inches(0.4),
                title, font_size=18, color=ACCENT_BLUE, bold=True)
    add_textbox(slide, Inches(1.1), top + Inches(0.45), Inches(11.5), Inches(1.3),
                desc, font_size=14, color=WHITE)


# ═══════════════════════════════════════════════════════════
# SLIDE 16 — SUGGESTED PATH / NEXT STEPS
# ═══════════════════════════════════════════════════════════
slide = dark_slide()
add_title_block(slide, "Suggested Development Path", "From benchtop to field deployment", 16, TOTAL_SLIDES)

phases = [
    ("Phase 1", "Benchtop Validation",
     "TI IWR6843ISK ($299)\nBuilt-in vital sign demo\nTest: respiration detection\nthrough fabric, plastic,\ncardboard at 0.5–2 m",
     ACCENT_BLUE),
    ("Phase 2", "Custom MIMO Array",
     "(If higher 3D resolution\nis required)\nCascade multiple IWR6843\nunits or engage Vayyar",
     ACCENT_GREEN),
    ("Phase 3", "ML Dataset Collection",
     "Labeled recordings of:\n• Live animals\n• Inert objects\n• Various concealment\n  scenarios",
     ACCENT_BLUE),
    ("Phase 4", "Field Trial",
     "Controlled checkpoint\nMeasure:\n• Pd (detection probability)\n• Pfa (false alarm rate)\nIterate & improve",
     ACCENT_GREEN),
]

for i, (label, title, desc, clr) in enumerate(phases):
    left = Inches(0.6 + i * 3.15)
    add_bar(slide, left, Inches(1.5), Inches(2.8), Inches(0.55), clr)
    add_textbox(slide, left + Inches(0.1), Inches(1.52), Inches(2.6), Inches(0.25),
                label, font_size=12, color=DARK_BG, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.1), Inches(1.78), Inches(2.6), Inches(0.25),
                title, font_size=14, color=DARK_BG, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.1), Inches(2.3), Inches(2.6), Inches(3.0),
                desc, font_size=13, color=WHITE, alignment=PP_ALIGN.CENTER)

    if i < len(phases) - 1:
        add_textbox(slide, left + Inches(2.82), Inches(2.6), Inches(0.4), Inches(0.4),
                    "▶", font_size=24, color=ACCENT_BLUE, bold=True)


# ── save ──────────────────────────────────────────────────
out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                        "mmWave_Radar_Live_Animal_Detection.pptx")
prs.save(out_path)
print(f"Saved to: {out_path}")
print(f"Slides: {len(prs.slides)}")
